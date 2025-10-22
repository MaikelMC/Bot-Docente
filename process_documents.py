import os
import pickle
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from config import DATA_DIR, VECTORSTORE_DIR, CHUNK_SIZE, CHUNK_OVERLAP, EMBEDDING_MODEL

def extract_from_pdf(filepath):
    """Extrae texto de PDF"""
    reader = PdfReader(filepath)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_from_docx(filepath):
    """Extrae texto de Word (.docx)"""
    doc = Document(filepath)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    
    # También extraer texto de tablas
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + " "
        text += "\n"
    
    return text

def extract_from_txt(filepath):
    """Extrae texto de archivo .txt"""
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def extract_from_xlsx(filepath):
    """Extrae texto de Excel (.xlsx)"""
    workbook = openpyxl.load_workbook(filepath, data_only=True)
    text = ""
    
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        text += f"\n=== Hoja: {sheet_name} ===\n"
        
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
            if row_text.strip():
                text += row_text + "\n"
    
    return text

def extract_from_pptx(filepath):
    """Extrae texto de PowerPoint (.pptx)"""
    prs = Presentation(filepath)
    text = ""
    
    for i, slide in enumerate(prs.slides, 1):
        text += f"\n=== Diapositiva {i} ===\n"
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    
    return text

def extract_from_csv(filepath):
    """Extrae texto de CSV"""
    import csv
    text = ""
    
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f)
        for row in reader:
            text += " | ".join(row) + "\n"
    
    return text

def extract_text_from_documents(doc_dir):
    """Extrae texto de todos los documentos soportados"""
    all_text = []
    
    # Formatos soportados
    supported_formats = {
        '.pdf': extract_from_pdf,
        '.docx': extract_from_docx,
        '.txt': extract_from_txt,
        '.xlsx': extract_from_xlsx,
        '.xls': extract_from_xlsx,
        '.pptx': extract_from_pptx,
        '.csv': extract_from_csv,
        '.md': extract_from_txt,  # Markdown como texto plano
        '.log': extract_from_txt,
        '.json': extract_from_txt,
        '.xml': extract_from_txt,
    }
    
    print("📄 Extrayendo texto de documentos...")
    print(f"   Formatos soportados: {', '.join(supported_formats.keys())}")
    
    files_processed = 0
    
    for filename in os.listdir(doc_dir):
        filepath = os.path.join(doc_dir, filename)
        
        # Verificar si es un archivo (no carpeta)
        if not os.path.isfile(filepath):
            continue
        
        # Obtener extensión
        file_ext = os.path.splitext(filename)[1].lower()
        
        if file_ext in supported_formats:
            print(f"  - Procesando: {filename}")
            
            try:
                # Llamar a la función correspondiente
                extractor = supported_formats[file_ext]
                text = extractor(filepath)
                
                if text and len(text.strip()) > 0:
                    all_text.append({
                        'filename': filename,
                        'text': text,
                        'format': file_ext
                    })
                    files_processed += 1
                else:
                    print(f"    ⚠️ Advertencia: {filename} está vacío")
                    
            except Exception as e:
                print(f"    ⚠️ Error procesando {filename}: {e}")
        else:
            print(f"  ⊘ Ignorado (formato no soportado): {filename}")
    
    print(f"\n  ✓ Procesados {files_processed} documentos exitosamente")
    return all_text

def split_text_into_chunks(documents, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    """Divide el texto en chunks con overlap"""
    chunks = []
    print("\n✂️ Dividiendo texto en chunks...")
    
    for doc in documents:
        text = doc['text']
        filename = doc['filename']
        file_format = doc.get('format', 'unknown')
        
        # Dividir en chunks
        for i in range(0, len(text), chunk_size - overlap):
            chunk = text[i:i + chunk_size]
            if len(chunk) > 50:  # Ignorar chunks muy pequeños
                chunks.append({
                    'text': chunk,
                    'source': filename,
                    'format': file_format,
                    'chunk_id': len(chunks)
                })
    
    print(f"  ✓ Creados {len(chunks)} chunks")
    return chunks

def create_embeddings(chunks):
    """Crea embeddings para cada chunk"""
    print("\n🧠 Generando embeddings (esto puede tardar)...")
    
    # Cargar modelo local
    model = SentenceTransformer(EMBEDDING_MODEL)
    
    # Extraer solo los textos
    texts = [chunk['text'] for chunk in chunks]
    
    # Generar embeddings
    embeddings = model.encode(texts, show_progress_bar=True)
    
    return embeddings, model

def create_faiss_index(embeddings):
    """Crea índice FAISS"""
    print("\n📊 Creando índice FAISS...")
    
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(embeddings.astype('float32'))
    
    print(f"  ✓ Índice creado con {index.ntotal} vectores")
    return index

def save_vectorstore(index, chunks):
    """Guarda el índice y los chunks"""
    print(f"\n💾 Guardando vectorstore en {VECTORSTORE_DIR}...")
    
    os.makedirs(VECTORSTORE_DIR, exist_ok=True)
    
    # Guardar índice FAISS
    faiss.write_index(index, os.path.join(VECTORSTORE_DIR, "index.faiss"))
    
    # Guardar chunks
    with open(os.path.join(VECTORSTORE_DIR, "chunks.pkl"), 'wb') as f:
        pickle.dump(chunks, f)
    
    print("  ✓ Vectorstore guardado exitosamente")

def main():
    """Función principal"""
    print("=" * 50)
    print("🚀 INICIANDO PROCESAMIENTO DE DOCUMENTOS")
    print("=" * 50)
    
    # 1. Verificar que exista el directorio
    if not os.path.exists(DATA_DIR):
        print(f"❌ Error: La carpeta {DATA_DIR} no existe")
        return
    
    # Contar archivos soportados
    supported_exts = ['.pdf', '.docx', '.txt', '.xlsx', '.xls', '.pptx', '.csv', '.md', '.log', '.json', '.xml']
    files = [f for f in os.listdir(DATA_DIR) 
            if os.path.isfile(os.path.join(DATA_DIR, f)) 
            and os.path.splitext(f)[1].lower() in supported_exts]
    
    if not files:
        print(f"❌ Error: No hay archivos soportados en {DATA_DIR}")
        print(f"   Formatos soportados: {', '.join(supported_exts)}")
        return
    
    print(f"✓ Encontrados {len(files)} archivos soportados")
    
    # 2. Extraer texto
    documents = extract_text_from_documents(DATA_DIR)
    
    if not documents:
        print("❌ Error: No se pudo extraer texto de ningún documento")
        return
    
    # 3. Dividir en chunks
    chunks = split_text_into_chunks(documents)
    
    if not chunks:
        print("❌ Error: No se generaron chunks")
        return
    
    # 4. Crear embeddings
    embeddings, model = create_embeddings(chunks)
    
    # 5. Crear índice FAISS
    index = create_faiss_index(embeddings)
    
    # 6. Guardar todo
    save_vectorstore(index, chunks)
    
    print("\n" + "=" * 50)
    print("✅ PROCESAMIENTO COMPLETADO")
    print("=" * 50)
    print(f"Documentos procesados: {len(documents)}")
    print(f"Total de chunks: {len(chunks)}")
    print(f"Dimensión de embeddings: {embeddings.shape[1]}")
    
    # Mostrar estadísticas por formato
    format_stats = {}
    for chunk in chunks:
        fmt = chunk['format']
        format_stats[fmt] = format_stats.get(fmt, 0) + 1
    
    print("\n📊 Chunks por formato:")
    for fmt, count in sorted(format_stats.items()):
        print(f"   {fmt}: {count} chunks")
    
if __name__ == "__main__":
    main()